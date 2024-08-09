import json
import os
import time

import pandas as pd
import plotly.io as pio
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from sklearn.model_selection import train_test_split

API_KEY = None
BASE_URL = "api.akkio.com/api"
URL = "api.akkio.com"
VERSION = "v1"
PROTOCOL = "https"
ENDPOINT = "chat-explore"
PORT = "443"


# Chat Wrappers
def create_chat_request(project_id, content):
    """
    Make API request for chat creation

    Returns:
        dict: json response
    """

    MODE = "new"

    url = f"{PROTOCOL}://{BASE_URL}/{VERSION}/{ENDPOINT}/{MODE}"
    headers = {"X-API-Key": API_KEY, "Content-Type": "application/json"}
    data = {
        "project_id": project_id,
        "messages": [{"role": "user", "content": content, "images": []}],
    }

    response = requests.post(
        url,
        json=data,
        headers=headers,
        timeout=120,
    )

    # Check for HTTP errors
    response.raise_for_status()

    return response.json()


def check_task_status(task_id):
    """
    Check task status from Chat creation POST call

    Returns:
        dict: json response
    """

    MODE = "status"

    url = f"{PROTOCOL}://{BASE_URL}/{VERSION}/{ENDPOINT}/{MODE}/{task_id}"
    headers = {"X-API-Key": API_KEY, "Content-Type": "application/json"}

    response = requests.get(
        url,
        headers=headers,
        timeout=120,
    )

    # Check for HTTP errors
    response.raise_for_status()

    return response.json()


def get_chat_results(chat_id, format_type="plotly_json"):
    """
    Get chat results based on chat_id

    Returns:
        dict: json response
    """

    # BASE64_PNG = "base64_png"
    # PLOTLY_JSON = "plotly_json"

    MODE = "chats"

    url = f"{PROTOCOL}://{BASE_URL}/{VERSION}/{ENDPOINT}/{MODE}/{chat_id}"
    headers = {"X-API-Key": API_KEY, "Content-Type": "application/json"}
    query_params = {"image_format": format_type}

    response = requests.get(
        url,
        headers=headers,
        params=query_params,
        timeout=120,
    )

    # Check for HTTP errors
    response.raise_for_status()
    return response.json()


def process_chat_output(data, output_file_path="ce"):

    raw_message = data["messages"][1]

    # Check if 'images' exist in raw_message
    if "images" in raw_message and raw_message["images"]:
        # Deserialize the first image's JSON data
        fig = pio.from_json(raw_message["images"][0])

        # Define desired dimensions and scale
        output_file_path = output_file_path + "_image.png"
        # fig.write_image(output_file_path, width=1920, height=1080, scale=1)

        fig.update_layout(
            xaxis=dict(
                tickangle=-45  # Slant the labels by 45 degrees. You can change this value as needed.
            )
        )

        fig.write_image(output_file_path, width=800, height=600, scale=2)

    # Check if 'table' exists in raw_message
    elif "table" in raw_message and raw_message["table"]:
        # Process table data here

        df = pd.DataFrame(raw_message["table"])

        output_file_path = output_file_path + "_table.csv"

        df.to_csv(output_file_path, index=False)

    # Fallback to process text if no images or table
    else:
        # Process text data here
        output_file_path = output_file_path + "_text.txt"
        with open(output_file_path, "w", encoding="utf-8") as file:
            file.write(raw_message["content"])


# Import data from disk
def import_data(filepath):
    """
    Reads a file (CSV, Excel) into a Pandas DataFrame.

    :param filepath: str, path to the file
    :return: DataFrame
    """

    # Mapping file extensions to Pandas functions
    file_readers = {
        ".csv": pd.read_csv,
        ".xlsx": pd.read_excel,
    }

    # Extract file extension and read the file
    file_extension = filepath.rsplit(".", 1)[-1].lower()
    read_function = file_readers.get("." + file_extension)

    if read_function:
        df = read_function(filepath, index_col=None)
    else:
        supported_types = ", ".join(file_readers.keys())
        err_msg = f"Unsupported file type '{file_extension}'. Supported file types are: {supported_types}."
        print(err_msg)
        raise ValueError(err_msg)

    return df


def df_to_dict(df):
    """
    Converts DataFrame to list of dictionaries
    """
    # Convert all data in the DataFrame to strings (force conversion to string since
    # we implicitely cast to the correct DType when generating the dataset view)
    df = df.applymap(str)

    # Convert DataFrame to a list of dictionaries
    dict_list = df.to_dict(orient="records")

    return dict_list


def add_rows_to_dataset(dataset_id, input_data):
    """
    Make API request to add rows to existing dataset

    Returns:
        dict: json response
    """
    start_time = time.time()
    response = requests.post(
        "{}://{}:{}/{}/datasets".format(PROTOCOL, URL, PORT, VERSION),
        json={"api_key": API_KEY, "id": dataset_id, "rows": input_data},
        timeout=120,
    )
    end_time = time.time()

    elapsed_time = end_time - start_time

    print(
        f"Request to add rows to dataset {dataset_id} with {len(input_data)} samples completed in {elapsed_time:.4f} seconds."
    )

    return response.json()


def partition_data(
    filename, out_location=os.getcwd(), test_partition=0.1, shuffle=True
):
    """Partition data to generate a train / test split"""

    fname = os.path.splitext(filename)[0]
    ext = os.path.splitext(filename)[1]

    df = pd.read_csv(filename)

    training_data, testing_data = train_test_split(
        df, test_size=test_partition, shuffle=shuffle, random_state=25
    )

    print(f"No. of training examples: {training_data.shape[0]}")
    print(f"No. of testing examples: {testing_data.shape[0]}")

    train_filename = os.path.join(out_location, fname + "-train" + ext)
    test_filename = os.path.join(out_location, fname + "-test" + ext)

    training_data.to_csv(train_filename, index=False)
    testing_data.to_csv(test_filename, index=False)


def create_project(project_name, owner_id, org_id):
    """
    Make API request for create project

    Returns:
        dict: json response
    """

    url = f"{PROTOCOL}://{BASE_URL}/{VERSION}/projects"
    headers = {"X-Api-Key": API_KEY, "Content-Type": "application/json"}
    data = {
        "name": project_name,
        "_owner": owner_id,
        "_org": org_id,
    }

    response = requests.post(
        url,
        json=data,
        headers=headers,
        timeout=120,
    )

    # Check for HTTP errors
    response.raise_for_status()

    return response.json()


def make_prediction(
    model_id,
    input_data,
    show_factors=False,
    save=False,
    save_file_path="",
):
    """
    Make API request for inference on new data

    Returns:
        dict: json response
    """
    start_time = time.time()
    response = requests.post(
        "{}://{}:{}/{}/models".format(PROTOCOL, URL, PORT, VERSION),
        json={
            "api_key": API_KEY,
            "sample": True,
            "id": model_id,
            "data": input_data,
            "show_factors": show_factors,
        },
        timeout=120,
    )
    end_time = time.time()

    elapsed_time = end_time - start_time

    print(
        f"Request to make predictions {model_id} with {len(input_data)} samples completed in {elapsed_time:.4f} seconds."
    )

    # Check for HTTP errors
    response.raise_for_status()

    # Parse JSON response
    resp_dict = response.json()

    # Check for application level errors
    if resp_dict.get("status") == "error":
        raise Exception(
            f"Error from API: {resp_dict.get('message', 'No error message provided')}"
        )

    if save:
        if not save_file_path:
            save_file_path = "predictions.csv"

        print(f"Saving prediction to disk to {save_file_path}")

        if "predictions" in resp_dict.keys():
            df = pd.DataFrame(resp_dict["predictions"])
        else:
            df = pd.DataFrame(resp_dict)

        df.to_csv(save_file_path, index=False)

        print("Done!")

    return resp_dict


def set_dataset_fields(dataset_id, fields):
    """
    Make API request to set dataset fields

    Returns:
        dict: json response
    """

    # filed column type options are below
    # - category
    # - id
    # - integer
    # - float
    # - string
    # - date
    # - unknown (maps to disabled)

    # input data must look like this
    # fields = [
    #     {
    #         'name': {column1_name},
    #         'type': {column1_type},
    #         'valid': true
    #     },
    #     {
    #         'name': {column2_name},
    #         'type': {column2_type},
    #         'valid': true
    #     },
    #     ...
    # ]

    start_time = time.time()
    response = requests.post(
        "{}://{}:{}/{}/datasets".format(PROTOCOL, URL, PORT, VERSION),
        json={"api_key": API_KEY, "id": dataset_id, "fields": fields},
        timeout=120,
    )
    end_time = time.time()

    elapsed_time = end_time - start_time

    print(
        f"Request to set dataset fields in {dataset_id} completed in {elapsed_time:.4f} seconds."
    )

    return response.json()


class PPTXExporter:
    """
    A class to export images and CSV data as slides in a PowerPoint presentation.

    Attributes:
    -----------
    artifacts_folder : str
        Path to the folder containing images and CSV files to be added to the presentation.
    prs : pptx.Presentation
        A PowerPoint presentation object used to create and manipulate slides.
    image_layout : dict
        A dictionary containing the layout configuration for images on slides.
        Default layout is {'left': Inches(1), 'top': Inches(1), 'width': Inches(5), 'height': Inches(5)}.
    table_layout : dict
        A dictionary containing the layout configuration for tables on slides.
        Default layout is {'left': Inches(1), 'top': Inches(1), 'width': Inches(8), 'height': Inches(5)}.

    Methods:
    --------
    create():
        Creates slides from images and CSV files found in the artifacts folder.
    save(filename):
        Saves the PowerPoint presentation to the specified filename.
    """

    def __init__(
        self, input_folder, image_layout=None, table_layout=None, text_layout=None
    ):
        """
        Initializes the PPTXExporter with the folder path containing artifacts and optional layout configurations.

        Parameters:
        -----------
        input_folder : str
            The folder path where images and CSV files are stored.
        image_layout : dict, optional
            Layout configuration for images (default is None, which uses a preset layout).
        table_layout : dict, optional
            Layout configuration for tables (default is None, which uses a preset layout).
        """
        self.artifacts_folder = input_folder
        self.prs = Presentation()

        self.slide_width = self.prs.slide_width
        self.slide_height = self.prs.slide_height

        self.table_layout = table_layout or {
            "left": Inches(1),
            "top": Inches(1),
            "width": Inches(8),
            "height": Inches(5),
        }

        self.text_layout = text_layout or {
            "left": Inches(1),
            "top": Inches(1),
            "width": Inches(8),
            "height": Inches(5),
        }

        # Adjust image layout to take up most of the slide, relative to slide size
        self.image_layout = image_layout or {
            "left": self.slide_width * 0.05,  # 5% margin from the left
            "top": self.slide_height * 0.05,  # 5% margin from the top
            "width": self.slide_width * 0.9,  # Width to cover 90% of the slide width
            "height": self.slide_height
            * 0.85,  # Height to cover 85% of the slide height
        }

    def create(self):
        """
        Creates slides for each image and CSV file in the artifacts folder.

        The method iterates through all files in the specified folder.
        If the file is an image (supported formats: PNG, JPG, JPEG, BMP, GIF, TIFF),
        it creates an image slide. If the file is a CSV file, it creates a table slide.
        """
        print("Starting process...")

        # Check if there are any files in the artifacts folder
        if not os.listdir(self.artifacts_folder):
            raise ValueError(
                f"No files found in {self.artifacts_folder}. Files need to be present."
            )

        for fname in os.listdir(self.artifacts_folder):
            print(f"Creating slide for: {fname}")
            file_path = os.path.join(self.artifacts_folder, fname)

            if fname.lower().endswith(
                (".png", ".jpg", ".jpeg", ".bmp", ".gif", ".tiff")
            ):
                self._add_image_slide(file_path)
            elif fname.lower().endswith(".csv"):
                self._add_table_slide(file_path)
            elif fname.lower().endswith(".txt"):
                self._add_text_slide(file_path)

    def _add_image_slide(self, image_path):
        """
        Adds an image slide to the presentation.

        Parameters:
        -----------
        image_path : str
            The full path to the image file to be added to the slide.
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        slide.shapes.add_picture(
            image_path,
            self.image_layout["left"],
            self.image_layout["top"],
            width=self.image_layout["width"],
            height=self.image_layout["height"],
        )

    def _add_table_slide(self, csv_file):
        """
        Adds a table slide to the presentation using data from a CSV file.

        Parameters:
        -----------
        csv_file : str
            The full path to the CSV file to be added to the slide.
        """
        df = pd.read_csv(csv_file)
        rows, cols = df.shape

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        table = slide.shapes.add_table(
            rows + 1,
            cols,
            self.table_layout["left"],
            self.table_layout["top"],
            self.table_layout["width"],
            self.table_layout["height"],
        ).table

        # Set column headers from the DataFrame
        for col_idx, col_name in enumerate(df.columns):
            table.cell(0, col_idx).text = col_name

        # Populate the table with CSV data
        for row_idx in range(rows):
            for col_idx in range(cols):
                table.cell(row_idx + 1, col_idx).text = str(df.iat[row_idx, col_idx])

    def _add_text_slide(self, text_file):
        """
        Adds a text slide to the presentation using content from a text file.

        Parameters:
        -----------
        text_file : str
            The full path to the text file to be added to the slide.
        """
        with open(text_file, "r", encoding="utf-8") as file:
            content = file.read()

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        textbox = slide.shapes.add_textbox(
            self.text_layout["left"],
            self.text_layout["top"],
            self.text_layout["width"],
            self.text_layout["height"],
        )
        text_frame = textbox.text_frame
        text_frame.text = content

        # Enable word wrap and adjust the text size if needed
        text_frame.word_wrap = True
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(12)  # Adjust the font size as needed

    def save(self, filename):
        """
        Saves the PowerPoint presentation to the specified file.

        Parameters:
        -----------
        filename : str
            The name of the file to save the presentation to. Must end with '.pptx'.

        Raises:
        -------
        ValueError
            If the filename does not end with '.pptx'.
        """
        if not filename.endswith(".pptx"):
            raise ValueError("filename must end with .pptx")

        print(f"Saving pptx to: {filename}")
        self.prs.save(filename)


if __name__ == "__main__":
    print("Starting process...")
